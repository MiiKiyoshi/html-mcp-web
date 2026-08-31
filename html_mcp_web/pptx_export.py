"""Turn a built slide deck into an editable pptx with firefox marionette and python-pptx.

Each block inside a body page's content becomes its own shape at the position it has in
the html (a 1280x720 px page is 12192000x6858000 EMU, 9525 EMU per px):
  - text blocks without KaTeX (p, headings, li, table-free inline text) -> text boxes
  - ul / ol                                                             -> bulleted / numbered text boxes
  - <img>                                                               -> the image file
  - tables without KaTeX                                                -> pptx tables
  - a div that paints a background or border                            -> a rounded panel, then its content
  - inline <svg> (self-contained)                                       -> real vector, with a screenshot fallback
  - KaTeX, canvas, anything else                                        -> a 3x screenshot of that element
The title bar, footer, corner logos, and page number bake into the slide background so
they cannot be grabbed; only the content stays editable over them. A page without a body
box (cover, contents, divider) bakes the whole page and overlays only its plain text.
Every skin is handled the same way; its look comes through the HTML it styles, so no skin
carries a pptx.

Skin configuration lives in skin.json under "pptx":
  fonts        the deck face, embedded so the pptx renders the same on a machine that
               lacks it: {"family": "Noto Sans KR", "regular": "fonts/X-Regular.ttf",
               "bold": ..., "italic": ..., "boldItalic": ...}, paths beside skin.json.
               Static TrueType (glyf) files whose fsType permits embedding; PowerPoint
               ignores CFF/OTF and variable fonts. Every run is set in family.
  font         typeface for every run when nothing is embedded; default Arial

Requires firefox, marionette_driver, python-pptx, and fontTools when fonts are embedded.
"""
import base64
import io
import json
import re
import shutil
import socket
import struct
import subprocess
import tempfile
import time
import zipfile
from pathlib import Path
from typing import Any

EMU_PER_PX = 9525
DPR = 3
PAGE_WIDTH_PX = 1280
PAGE_HEIGHT_PX = 720
PT_PER_PX = 0.75
DEFAULT_FONT = "Arial"
MONO_FONT = "Consolas"
BULLETS = ["•", "◦", "▪"]

JS_SETUP = """
const st = document.createElement('style');
st.textContent = 'section.page{width:%dpx !important;height:%dpx !important;margin:0 auto !important;'
  + 'box-shadow:none !important;border-radius:0 !important}.script-block{display:none !important}';
document.head.appendChild(st);
""" % (PAGE_WIDTH_PX, PAGE_HEIGHT_PX)

JS_READY = ("return document.readyState === 'complete' && "
            "Array.from(document.images).every(im => im.complete)")

JS_PAGES = "return document.querySelectorAll('section.page').length;"

# Hide the marked elements on one page so the background screenshot leaves their place empty
# for the editable shapes laid over it. A ::before decoration would go with the element, so
# the full-bleed pass hides only plain text; a body page hides its own content blocks.
JS_HIDE = ("const page = document.querySelectorAll('section.page')[arguments[0]];"
           "for (const i of arguments[1]) { const el = page.querySelector('[data-pptx-index=\"' + i + '\"]');"
           "if (el) el.style.visibility = 'hidden'; }")

# Walk a page's DOM into a flat list of shapes. Containers descend; a container that
# paints something becomes a panel first so its content sits on it.
JS_EXTRACT = r"""
const idx = arguments[0];
const page = document.querySelectorAll('section.page')[idx];
const pr = page.getBoundingClientRect();
const cs = (el) => getComputedStyle(el);
const rel = (r) => [r.left - pr.left, r.top - pr.top, r.width, r.height];
const BLOCKISH = /^(block|flex|grid|table|list-item|flow-root)$/;
const painted = (s) => (s.backgroundColor !== 'rgba(0, 0, 0, 0)' && s.backgroundColor !== 'transparent')
  || (parseFloat(s.borderTopWidth) > 0 && s.borderTopStyle !== 'none')
  || (parseFloat(s.borderLeftWidth) > 0 && s.borderLeftStyle !== 'none')
  || (parseFloat(s.borderBottomWidth) > 0 && s.borderBottomStyle !== 'none')
  || (parseFloat(s.borderRightWidth) > 0 && s.borderRightStyle !== 'none');
const items = [];
let counter = 0;
for (const stale of page.querySelectorAll('[data-pptx-index]')) stale.removeAttribute('data-pptx-index');
const mark = (el) => { el.setAttribute('data-pptx-index', String(counter)); return counter++; };

// A numbered marker: a round fill with a digit on it. A run carries no shape of its own,
// and both ways round it fail. A highlight is a rectangle that some readers carry on past
// a line break, and a shape placed at the marker's own coordinates drifts from the digit,
// because the text is laid out again in the deck's own metrics. The character that already
// means this is used instead: it sits in the line, so it cannot drift.
const luminance = (colour) => {
  const parts = (colour.match(/[\d.]+/g) || []).slice(0, 3).map(Number);
  return parts.length < 3 ? 1 : (0.2126 * parts[0] + 0.7152 * parts[1] + 0.0722 * parts[2]) / 255;
};
const badgeRun = (el, s, base) => {
  const rects = el.getClientRects();
  if (rects.length !== 1) return null;
  const radius = /%$/.test(s.borderTopLeftRadius)
    ? parseFloat(s.borderTopLeftRadius) / 100
    : (parseFloat(s.borderTopLeftRadius) || 0) / Math.max(1, Math.min(rects[0].width, rects[0].height));
  if (radius < 0.4) return null;
  const number = Number(el.textContent.trim());
  if (!Number.isInteger(number) || number < 1 || number > 10) return null;
  // Filled when the disc is darker than what is written on it, hollow when it is not; the
  // glyph is one colour and the numeral is left open, which is how the page draws it too.
  const filled = luminance(s.backgroundColor) < luminance(s.color);
  return Object.assign({}, base, {
    text: String.fromCodePoint((filled ? 0x2776 : 0x2460) + number - 1),
    color: filled ? s.backgroundColor : s.color,
    size: parseFloat(s.fontSize), mono: false, chip: null,
  });
};

// A code block keeps its line breaks and its indentation: the text is the layout there,
// and collapsing it ran a listing together into one line.
const keepsBreaks = (el) => el !== null && /^(pre|pre-wrap|break-spaces)/.test(cs(el).whiteSpace);

const runsOf = (node, out, base) => {
  for (const n of node.childNodes) {
    if (n.nodeType === 3) {
      if (keepsBreaks(n.parentElement)) {
        const lines = n.textContent.split('\n');
        for (let index = 0; index < lines.length; index++) {
          if (index > 0 && out.length > 0) out.push(Object.assign({}, base, {text: '\v'}));
          if (lines[index] !== '') out.push(Object.assign({}, base, {text: lines[index]}));
        }
        continue;
      }
      const t = n.textContent.replace(/\s+/g, ' ');
      if (t.trim() === '' && out.length === 0) continue;
      if (t !== '') out.push(Object.assign({}, base, {text: t}));
    } else if (n.nodeType === 1) {
      const tg = n.tagName.toLowerCase();
      if (tg === 'br') { out.push(Object.assign({}, base, {text: '\v'})); continue; }
      if (tg === 'ul' || tg === 'ol') continue;
      const s = cs(n);
      if (s.display === 'none') continue;
      const mono = /mono/i.test(s.fontFamily) || tg === 'code' || base.mono;
      const painted = s.backgroundColor !== 'rgba(0, 0, 0, 0)' && s.backgroundColor !== 'transparent';
      const marker = painted ? badgeRun(n, s, base) : null;
      if (marker) { out.push(marker); continue; }
      const chipBg = (painted && tg === 'code') ? s.backgroundColor : base.chip;
      const b = Object.assign({}, base, {
        bold: (parseInt(s.fontWeight) >= 600) || base.bold,
        italic: s.fontStyle === 'italic' || base.italic,
        color: s.color, size: parseFloat(s.fontSize), mono,
        spacing: s.letterSpacing === 'normal' ? 0 : (parseFloat(s.letterSpacing) || 0),
        chip: chipBg,
        link: tg === 'a' ? n.getAttribute('href') : base.link,
        sup: tg === 'sup' || base.sup, sub: tg === 'sub' || base.sub,
      });
      runsOf(n, out, b);
    }
  }
};
const paraOf = (el, level, marginTop) => {
  const s = cs(el);
  const runs = [];
  runsOf(el, runs, {bold: parseInt(s.fontWeight) >= 600, italic: s.fontStyle === 'italic',
                    color: s.color, size: parseFloat(s.fontSize), mono: /mono/i.test(s.fontFamily),
                    spacing: s.letterSpacing === 'normal' ? 0 : (parseFloat(s.letterSpacing) || 0), chip: null,
                    link: null, sup: false, sub: false});
  while (runs.length && runs[0].text.trim() === '' && runs[0].text !== '\v') runs.shift();
  while (runs.length && runs[runs.length - 1].text.trim() === '' && runs[runs.length - 1].text !== '\v') runs.pop();
  return {level, runs, size: parseFloat(s.fontSize),
          lh: parseFloat(s.lineHeight) / parseFloat(s.fontSize) || 1.2,
          align: s.textAlign, marginTop: marginTop || 0,
          // The gap under a paragraph is as much a part of the page as the gap over it, and
          // bullets are usually spaced with this one alone.
          marginBottom: parseFloat(s.marginBottom) || 0};
};
const listParas = (list, level, out) => {
  for (const li of list.children) {
    if (li.tagName.toLowerCase() !== 'li') continue;
    out.push(paraOf(li, level, parseFloat(cs(li).marginTop) || 0));
    for (const sub of li.children) {
      const t = sub.tagName.toLowerCase();
      if (t === 'ul' || t === 'ol') listParas(sub, level + 1, out);
    }
  }
};
// Positioned children (a footer's page number and label) keep their own boxes too, and so
// does a figure: an inline <svg> is not a block, so a box holding one read as a leaf and
// the whole drawing came out as a line of text made of its own labels.
const FIGURE = /^(svg|canvas|video|img)$/;
const hasBlockChild = (el) => Array.from(el.children).some((ch) => {
  const c = cs(ch);
  return BLOCKISH.test(c.display) || c.position === 'absolute' || c.position === 'fixed'
    || FIGURE.test(ch.tagName.toLowerCase());
});
// A ::before/::after that draws something (a contents ring number). Text carrying such a
// decoration cannot be hidden for the background pass without losing it, so it stays baked.
const pseudoContent = (n) => ['::before', '::after'].some((p) => {
  const c = cs(n, p).content; return c && c !== 'none' && c !== 'normal' && c !== '""' && c !== "''";
});
const hasPseudo = (el) => pseudoContent(el) || Array.from(el.querySelectorAll('*')).some(pseudoContent);

// Images in inline flow (a row of <img> in a div, default display:inline) are not block
// children, so their container reads as a leaf whose text is taken while the images are
// dropped. Each such image still needs its own picture frame.
const inlineImages = (el) => {
  for (const im of el.querySelectorAll('img')) {
    const ir = im.getBoundingClientRect();
    if (ir.width === 0 || ir.height === 0) continue;
    items.push({kind: 'img', rect: rel(ir), src: im.getAttribute('src'), i: mark(im)});
  }
};

const emit = (el) => {
  const tag = el.tagName.toLowerCase();
  if (tag === 'aside' || tag === 'script' || tag === 'style') return;
  const s = cs(el);
  if (s.display === 'none' || s.visibility === 'hidden') return;
  const r = el.getBoundingClientRect();
  if (r.width === 0 || r.height === 0) return;
  const rect = rel(r);
  const hasKatex = !!el.querySelector('.katex') || el.classList.contains('katex');
  const container = hasBlockChild(el);
  if (tag === 'img') {
    // Where the picture is actually painted: object-fit decides how it sits in a box of
    // another shape, and taking the box alone stretched a tall figure across a wide one.
    items.push({kind: 'img', rect, src: el.getAttribute('src'), fit: s.objectFit,
                natural: [el.naturalWidth, el.naturalHeight], i: mark(el)});
    return;
  }
  // A block that holds math (a paragraph, a list, a table) is shot whole; a container
  // that merely has math somewhere inside descends so the rest stays editable.
  const leafBlock = !container || tag === 'table' || tag === 'ul' || tag === 'ol';
  if (tag === 'svg' || tag === 'canvas' || tag === 'video' || (hasKatex && leafBlock)) {
    // A self-contained inline <svg> (styles on its own tags, no external CSS) travels as
    // real vector so it stays sharp; the screenshot still rides along as the fallback for
    // viewers older than PowerPoint 2016. KaTeX is HTML, not one svg, so it stays raster.
    const item = {kind: 'shot', rect, i: mark(el)};
    if (tag === 'svg') item.svg = new XMLSerializer().serializeToString(el);
    items.push(item); return;
  }
  if (tag === 'table') {
    const rows = [];
    for (const tr of el.querySelectorAll('tr')) {
      const cells = [];
      for (const td of tr.children) {
        const s2 = cs(td);
        // A cell that spans says so, and the row carries its own height: a spanning cell is
        // as tall as the rows it covers, which is not the height of the row it starts in.
        cells.push({rect: rel(td.getBoundingClientRect()), header: td.tagName.toLowerCase() === 'th',
                    colSpan: td.colSpan || 1, rowSpan: td.rowSpan || 1,
                    align: s2.textAlign, bg: s2.backgroundColor, para: paraOf(td, 0, 0),
                    borderBottom: [s2.borderBottomColor, parseFloat(s2.borderBottomWidth) || 0, s2.borderBottomStyle],
                    borderTop: [s2.borderTopColor, parseFloat(s2.borderTopWidth) || 0, s2.borderTopStyle],
                    borderLeft: [s2.borderLeftColor, parseFloat(s2.borderLeftWidth) || 0, s2.borderLeftStyle],
                    borderRight: [s2.borderRightColor, parseFloat(s2.borderRightWidth) || 0, s2.borderRightStyle]});
      }
      if (cells.length) rows.push({rect: rel(tr.getBoundingClientRect()), cells});
    }
    items.push({kind: 'table', rect, rows, i: mark(el)});
    return;
  }
  if (tag === 'ul' || tag === 'ol') {
    const paras = [];
    listParas(el, 0, paras);
    items.push({kind: 'text', rect, paras, list: tag, padLeft: parseFloat(s.paddingLeft) || 0, pseudo: hasPseudo(el), i: mark(el)});
    return;
  }
  if (painted(s)) {
    const side = (name) => [s['border' + name + 'Color'], s['border' + name + 'Style'] === 'none' ? 0 : (parseFloat(s['border' + name + 'Width']) || 0)];
    items.push({kind: 'panel', rect, bg: s.backgroundColor, radius: parseFloat(s.borderTopLeftRadius) || 0,
                borders: {top: side('Top'), right: side('Right'), bottom: side('Bottom'), left: side('Left')}, i: mark(el)});
    if (!container) {
      inlineImages(el);
      const para = paraOf(el, 0, 0);
      if (para.runs.length) items.push({kind: 'text', rect: [rect[0] + (parseFloat(s.paddingLeft) || 0), rect[1] + (parseFloat(s.paddingTop) || 0),
        rect[2] - (parseFloat(s.paddingLeft) || 0) - (parseFloat(s.paddingRight) || 0), rect[3] - (parseFloat(s.paddingTop) || 0) - (parseFloat(s.paddingBottom) || 0)],
        paras: [para], list: null, pseudo: hasPseudo(el), i: mark(el)});
      return;
    }
  }
  if (container) { for (const ch of el.children) emit(ch); return; }
  inlineImages(el);
  const para = paraOf(el, 0, 0);
  // plain: pure text on this element (no border/background/list, reached without a paint),
  // so it can be hidden for the background pass and re-added as an editable box with nothing
  // of its own lost. Painted or listed text (the other pushes) stays baked into the picture.
  if (para.runs.length) items.push({kind: 'text', rect, paras: [para], list: null, pseudo: hasPseudo(el), plain: true, i: mark(el)});
};

const body = page.querySelector('.body');
const titleEl = page.querySelector('.tbar h2');
const title = (titleEl || {}).textContent || '';
// The chrome (title bar, footer, corner logos, page number) bakes into the slide
// background so it cannot be grabbed; only the body's own blocks, tagged here, stay
// as editable shapes laid over it. The bar's title is the exception: the deck is edited
// per page by its title, so the words come back as a text box while the bar stays baked.
for (const ch of page.children) {
  if (ch === body || ch.classList.contains('script-block')) continue;
  emit(ch);
}
if (titleEl) {
  const marked = titleEl.getAttribute('data-pptx-index');
  if (marked !== null) {
    const item = items.find((it) => String(it.i) === marked && it.kind === 'text');
    if (item) item.body = true;
  }
}
const chromeCount = items.length;
if (body) for (const ch of body.children) emit(ch);
for (let k = chromeCount; k < items.length; k++) items[k].body = true;
return {title: title.trim(), hasBody: !!body, items};
"""

# Marks restart at 0 on every page, so the lookup stays inside the page being exported;
# a document-wide search returned an earlier page's element and its picture was
# stretched into this page's box.
JS_ELEMENT = """
const page = document.querySelectorAll('section.page')[arguments[0]];
const el = page.querySelector('[data-pptx-index="' + arguments[1] + '"]');
el.scrollIntoView({block: 'center'});
return el;
"""


def _px(value: float):
    from pptx.util import Emu
    return Emu(int(round(value * EMU_PER_PX)))


def _rgb(css: str | None):
    from pptx.dml.color import RGBColor
    match = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)", css or "")
    if not match or (match[4] is not None and float(match[4]) == 0):
        return None
    return RGBColor(int(match[1]), int(match[2]), int(match[3]))


def _set_font(run, name: str) -> None:
    from pptx.oxml.ns import qn
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        element = rPr.find(qn(tag))
        if element is None:
            element = rPr.makeelement(qn(tag), {})
            rPr.append(element)
        element.set("typeface", name)


def _clear_bullet(pPr) -> None:
    from pptx.oxml.ns import qn
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum", "a:buFont"):
        for element in pPr.findall(qn(tag)):
            pPr.remove(element)


def _set_list_marker(p_xml, level: int, pad_left: float, ordered: bool, font: str) -> None:
    from pptx.oxml.ns import qn
    pPr = p_xml.get_or_add_pPr()
    indent_px = pad_left + 18 * level
    pPr.set("marL", str(int(indent_px * EMU_PER_PX)))
    pPr.set("indent", str(int(-14 * EMU_PER_PX)))
    _clear_bullet(pPr)
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": font}))
    if ordered:
        pPr.append(pPr.makeelement(qn("a:buAutoNum"), {"type": "arabicPeriod"}))
    else:
        pPr.append(pPr.makeelement(qn("a:buChar"), {"char": BULLETS[min(level, 2)]}))


def _no_marker(p_xml) -> None:
    from pptx.oxml.ns import qn
    pPr = p_xml.get_or_add_pPr()
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    _clear_bullet(pPr)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _set_run_highlight(run, color) -> None:
    """A `<code>` chip's background becomes the run's highlight, the only inline background
    a pptx run has. It goes right after the fill so the rPr child order stays valid."""
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    for existing in rPr.findall(qn("a:highlight")):
        rPr.remove(existing)
    highlight = rPr.makeelement(qn("a:highlight"), {})
    srgb = highlight.makeelement(qn("a:srgbClr"), {"val": str(color)})
    highlight.append(srgb)
    fill = rPr.find(qn("a:solidFill"))
    if fill is not None:
        fill.addnext(highlight)
    else:
        rPr.insert(0, highlight)


def _fill_paragraph(p, para: dict[str, Any], font: str, bold_all: bool = False,
                    space_before: float | None = None) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    # Fixed line pitch in points: a percentage is taken against PowerPoint's own single
    # spacing, which is not the browser's, and the lines drift apart.
    p.line_spacing = Pt(para["lh"] * para["size"] * PT_PER_PX)
    gap = para["marginTop"] if space_before is None else space_before
    p.space_before = Pt(gap * PT_PER_PX) if gap else Pt(0)
    p.space_after = Pt(0)
    p.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(para["align"], PP_ALIGN.LEFT)
    if not para["runs"]:
        run = p.add_run()
        run.text = ""
        run.font.size = Pt(para["size"] * PT_PER_PX)
        _set_font(run, font)
        return
    for piece in para["runs"]:
        if piece["text"] == "\v":
            p.add_line_break()
            continue
        run = p.add_run()
        run.text = piece["text"]
        _set_font(run, MONO_FONT if piece["mono"] else font)
        size = piece["size"] * PT_PER_PX
        if piece["sup"] or piece["sub"]:
            size *= 0.75
            run._r.get_or_add_rPr().set("baseline", "30000" if piece["sup"] else "-25000")
        run.font.size = Pt(size)
        # Kern at every size, the way the browser does; PowerPoint otherwise leaves small
        # text unkerned and the Latin words read looser than on the page.
        run._r.get_or_add_rPr().set("kern", "0")
        run.font.bold = bool(piece["bold"]) or bold_all
        run.font.italic = bool(piece["italic"])
        color = _rgb(piece["color"])
        if color is not None:
            run.font.color.rgb = color
        # Tracking: the run's spc is in 1/100 pt, from the DOM's px letter-spacing.
        spacing = piece.get("spacing") or 0
        if spacing:
            run._r.get_or_add_rPr().set("spc", str(int(round(spacing * PT_PER_PX * 100))))
        chip = _rgb(piece.get("chip")) if piece.get("chip") else None
        if chip is not None:
            _set_run_highlight(run, chip)
        if piece["link"]:
            run.hyperlink.address = piece["link"]


def _add_text(slide, item: dict[str, Any], font: str) -> None:
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    x, y, w, h = item["rect"]
    box = slide.shapes.add_textbox(_px(x), _px(y), _px(w + 6), _px(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    # Centre the text in its box: a multi-line block fills its box so this is a no-op, while
    # a single line in a line-height-taller box (a bar title) is centred instead of riding
    # its top with empty space below.
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    bodyPr = frame._txBody.find(qn("a:bodyPr"))
    for element in list(bodyPr):
        if element.tag in (qn("a:spAutoFit"), qn("a:normAutofit")):
            bodyPr.remove(element)
    previous_bottom = 0.0
    for index, para in enumerate(item["paras"]):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if item["list"]:
            _set_list_marker(p._p, para["level"], item["padLeft"], item["list"] == "ol", font)
        else:
            _no_marker(p._p)
        # Neighbouring margins collapse on the page, where the gap is the larger of the two;
        # in a text box they would add, so the collapse is done here and carried above the
        # paragraph. Spacing bullets with margin-bottom alone is the common case, and taking
        # only margin-top closed every gap in the list.
        gap = para["marginTop"] if index == 0 else max(previous_bottom, para["marginTop"])
        _fill_paragraph(p, para, font, space_before=gap)
        previous_bottom = para.get("marginBottom", 0)


def _wider(current, other):
    """Of the two lines a pair of cells asks for along their shared edge, the one drawn."""
    if other is None or other[1] <= 0 or other[2] == "none":
        return current
    if current is None or other[1] > current[1]:
        return other
    return current


def _set_cell_borders(cell, top, bottom, left=None, right=None) -> None:
    """Each side: (css color, width px, style) from the DOM, or None for no line."""
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    for tag, spec in (("a:lnL", left), ("a:lnR", right), ("a:lnT", top), ("a:lnB", bottom)):
        for element in tcPr.findall(qn(tag)):
            tcPr.remove(element)
        line = tcPr.makeelement(qn(tag), {})
        color = _rgb(spec[0]) if spec and spec[1] > 0 and spec[2] != "none" else None
        if color is None:
            line.set("w", "0")
            line.append(line.makeelement(qn("a:noFill"), {}))
        else:
            line.set("w", str(int(spec[1] * 12700)))
            fill = line.makeelement(qn("a:solidFill"), {})
            fill.append(fill.makeelement(qn("a:srgbClr"), {"val": str(color)}))
            line.append(fill)
        tcPr.append(line)


def _add_table(slide, item: dict[str, Any], font: str) -> None:
    from pptx.enum.text import MSO_ANCHOR
    from pptx.oxml.ns import qn
    from pptx.util import Emu
    rows = item["rows"]
    n_rows = len(rows)
    n_cols = max(sum(cell["colSpan"] for cell in row["cells"]) for row in rows)
    x, y, w, h = item["rect"]
    graphic = slide.shapes.add_table(n_rows, n_cols, _px(x), _px(y), _px(w), _px(h))
    table = graphic.table
    tblPr = table._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")
    style = tblPr.find(qn("a:tableStyleId"))
    if style is not None:
        # No Style, No Grid: the style adds nothing, so the per-cell borders read from the
        # DOM are the only lines. The Table Grid style drew a full grid over them in PowerPoint.
        style.text = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    # A row holds one cell per slot it still has free, not one per column: a cell above it
    # spanning down, or one before it spanning across, has already taken the rest. Filling by
    # position in the row put a header group over the wrong columns and left the far side of
    # the row empty. Each slot records the cell it belongs to and whether it starts there.
    owner = [[None] * n_cols for _ in range(n_rows)]
    widths = [None] * n_cols
    placed = []
    for row_index, row in enumerate(rows):
        col = 0
        for source in row["cells"]:
            while col < n_cols and owner[row_index][col] is not None:
                col += 1
            if col >= n_cols:
                break
            reach_rows = min(source["rowSpan"], n_rows - row_index)
            reach_cols = min(source["colSpan"], n_cols - col)
            for down in range(reach_rows):
                for across in range(reach_cols):
                    owner[row_index + down][col + across] = (source, down == 0 and across == 0)
            if source["colSpan"] == 1:
                widths[col] = source["rect"][2]
            placed.append((row_index, col, reach_rows, reach_cols, source))
            col += source["colSpan"]
    # Two cells meet along one edge and the page draws a single line there. Writing each cell
    # its own four sides left one of a pair asking for a line where the other asked for none,
    # and the slots a spanning cell covers asked for nothing at all, so a reader that draws
    # each slot of the grid in turn broke the line into pieces. The edges are settled once,
    # here, and every slot reads its four sides back off them.
    horizontal, vertical = {}, {}
    for row_index, col, reach_rows, reach_cols, source in placed:
        for across in range(reach_cols):
            above, below = (row_index, col + across), (row_index + reach_rows, col + across)
            horizontal[above] = _wider(horizontal.get(above), source["borderTop"])
            horizontal[below] = _wider(horizontal.get(below), source["borderBottom"])
        for down in range(reach_rows):
            before, after = (row_index + down, col), (row_index + down, col + reach_cols)
            vertical[before] = _wider(vertical.get(before), source["borderLeft"])
            vertical[after] = _wider(vertical.get(after), source["borderRight"])
    spare = max(0.0, w - sum(width for width in widths if width is not None))
    unknown = sum(1 for width in widths if width is None)
    for col in range(n_cols):
        table.columns[col].width = _px(widths[col] if widths[col] is not None else spare / unknown)
    for row_index, row in enumerate(rows):
        table.rows[row_index].height = _px(row["rect"][3])
    reach = {(row_index, col): (reach_rows, reach_cols)
             for row_index, col, reach_rows, reach_cols, _ in placed}
    for row_index, col, reach_rows, reach_cols, source in placed:
        if reach_rows > 1 or reach_cols > 1:
            table.cell(row_index, col).merge(
                table.cell(row_index + reach_rows - 1, col + reach_cols - 1))
    for row_index in range(n_rows):
        for col in range(n_cols):
            slot = owner[row_index][col]
            source = None if slot is None else slot[0]
            starts = slot is not None and slot[1]
            cell = table.cell(row_index, col)
            cell.margin_left = cell.margin_right = Emu(int(6 * EMU_PER_PX))
            cell.margin_top = cell.margin_bottom = Emu(int(4 * EMU_PER_PX))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # A cell that starts a span is drawn over all of it, so its far sides are the
            # span's; a slot it covers is drawn on its own and takes the sides it sits on.
            down, across = reach[(row_index, col)] if starts else (1, 1)
            # Borders go in before the fill: the schema orders lnL/R/T/B ahead of the fill,
            # and PowerPoint drops borders that come after it (LibreOffice tolerates them).
            _set_cell_borders(cell,
                              horizontal.get((row_index, col)),
                              horizontal.get((row_index + down, col)),
                              vertical.get((row_index, col)),
                              vertical.get((row_index, col + across)))
            background = _rgb(source["bg"]) if source is not None else None
            if background is None:
                cell.fill.background()
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = background
            if not starts:
                continue
            frame = cell.text_frame
            frame.word_wrap = True
            _fill_paragraph(frame.paragraphs[0], source["para"], font, bold_all=source["header"])


def _add_panel(slide, item: dict[str, Any]) -> None:
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Pt
    x, y, w, h = item["rect"]
    borders = item["borders"]
    sides = {name: (_rgb(color), width) for name, (color, width) in borders.items()}
    drawn = {name: value for name, value in sides.items() if value[0] is not None and value[1] > 0}
    uniform = len(drawn) == 4 and len({value for value in drawn.values()}) == 1
    background = _rgb(item["bg"])
    if background is not None or uniform:
        kind = MSO_SHAPE.ROUNDED_RECTANGLE if item["radius"] > 0 else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(kind, _px(x), _px(y), _px(w), _px(h))
        if item["radius"] > 0:
            shape.adjustments[0] = min(0.5, item["radius"] / min(w, h))
        shape.shadow.inherit = False
        if background is None:
            shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb = background
        if uniform:
            color, width = next(iter(drawn.values()))
            shape.line.color.rgb = color
            shape.line.width = Pt(width * PT_PER_PX)
        else:
            shape.line.fill.background()
    if uniform:
        return
    # Partial borders (a rule under a header row, an accent bar on the left) are drawn as
    # thin rectangles on the sides that have them.
    for name, (color, width) in drawn.items():
        if name == "top":
            box = (x, y, w, width)
        elif name == "bottom":
            box = (x, y + h - width, w, width)
        elif name == "left":
            box = (x, y, width, h)
        else:
            box = (x + w - width, y, width, h)
        rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _px(box[0]), _px(box[1]), _px(box[2]), _px(box[3]))
        rule.shadow.inherit = False
        rule.fill.solid()
        rule.fill.fore_color.rgb = color
        rule.line.fill.background()


def _add_picture(slide, data: bytes, rect: list[float]):
    x, y, w, h = rect
    return slide.shapes.add_picture(io.BytesIO(data), _px(x), _px(y), _px(w), _px(h))


def _painted_rect(rect: list[float], natural: list[float] | None, fit: str | None) -> list[float]:
    """Where inside its box the picture is actually painted.

    A box of another shape than the picture does not stretch it unless the page says so:
    contain fits it inside and centres it, leaving a band on two sides. Reading the box
    alone put a tall figure across a wide one.
    """
    x, y, w, h = rect
    if not natural or fit not in ("contain", "scale-down") or w <= 0 or h <= 0:
        return rect
    natural_width, natural_height = natural
    if not natural_width or not natural_height:
        return rect
    scale = min(w / natural_width, h / natural_height)
    if fit == "scale-down":
        scale = min(scale, 1.0)
    width = natural_width * scale
    height = natural_height * scale
    # object-position defaults to the centre, which is what every figure in the templates uses.
    return [x + (w - width) / 2, y + (h - height) / 2, width, height]


def _crop_to_cover(picture, rect: list[float], natural: list[float] | None) -> None:
    """cover fills the box and the rest is cut off, so the picture carries the same crop."""
    if not natural:
        return
    natural_width, natural_height = natural
    if not natural_width or not natural_height:
        return
    _, _, w, h = rect
    scale = max(w / natural_width, h / natural_height)
    width = natural_width * scale
    height = natural_height * scale
    if width > w:
        picture.crop_left = picture.crop_right = (width - w) / width / 2
    if height > h:
        picture.crop_top = picture.crop_bottom = (height - h) / height / 2


def _set_slide_background(slide, png: bytes) -> None:
    """Put a full-page screenshot in as the slide's background fill, not a floating picture,
    so it fills the slide and the mouse cannot grab or move it."""
    from lxml import etree
    from pptx.oxml.ns import qn

    _, rid = slide.part.get_or_add_image_part(io.BytesIO(png))
    c_sld = slide._element.find(qn("p:cSld"))
    bg = c_sld.makeelement(qn("p:bg"), {})
    bg_pr = etree.SubElement(bg, qn("p:bgPr"))
    blip_fill = etree.SubElement(bg_pr, qn("a:blipFill"))
    etree.SubElement(blip_fill, qn("a:blip")).set(qn("r:embed"), rid)
    etree.SubElement(etree.SubElement(blip_fill, qn("a:stretch")), qn("a:fillRect"))
    etree.SubElement(bg_pr, qn("a:effectLst"))
    c_sld.insert(0, bg)


# PowerPoint 2016+ keeps a picture vector when the blip references an SVG through this
# extension; the blip's own r:embed stays the raster fallback for older viewers.
SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _normalize_svg(svg: bytes, width_px: float, height_px: float, font: str) -> bytes:
    """Match the SVG's viewBox aspect to the picture frame so PowerPoint, which fills the
    frame, does not stretch the drawing. The browser letterboxes (preserveAspectRatio meet);
    the viewBox is widened or heightened symmetrically to carry that same empty margin, so
    the drawing keeps its shape and stays centred. The deck font is named on the root so its
    <text> keeps the browser's metrics: an inline <svg> inherits the page font through CSS,
    but a standalone one carries no CSS, and PowerPoint would substitute a wider face and push
    labels over the drawing."""
    from lxml import etree

    root = etree.fromstring(svg)
    if root.get("font-family") is None:
        root.set("font-family", font)
    view_box = root.get("viewBox") or root.get("viewbox")
    if view_box:
        min_x, min_y, vb_w, vb_h = (float(value) for value in view_box.replace(",", " ").split())
        if vb_w > 0 and vb_h > 0 and width_px > 0 and height_px > 0:
            frame = width_px / height_px
            drawing = vb_w / vb_h
            if drawing < frame:
                widened = vb_h * frame
                min_x -= (widened - vb_w) / 2
                vb_w = widened
            elif drawing > frame:
                heightened = vb_w / frame
                min_y -= (heightened - vb_h) / 2
                vb_h = heightened
            root.set("viewBox", f"{min_x:.3f} {min_y:.3f} {vb_w:.3f} {vb_h:.3f}")
    root.attrib.pop("viewbox", None)
    root.set("preserveAspectRatio", "xMidYMid meet")
    root.set("width", f"{width_px:.2f}")
    root.set("height", f"{height_px:.2f}")
    root.attrib.pop("style", None)  # width:100%/max-height would fight the frame size
    return etree.tostring(root)


def _add_svg_picture(slide, png: bytes, svg: bytes, rect: list[float]) -> None:
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.oxml.ns import qn

    picture = _add_picture(slide, png, rect)
    package = slide.part.package
    part = Part(package.next_partname("/ppt/media/image%d.svg"), "image/svg+xml", package, svg)
    rid = slide.part.relate_to(part, RT.IMAGE)
    blip = picture._element.blipFill.blip
    ext_list = blip.find(qn("a:extLst"))
    if ext_list is None:
        ext_list = etree.SubElement(blip, qn("a:extLst"))
    ext = etree.SubElement(ext_list, qn("a:ext"))
    ext.set("uri", SVG_EXT_URI)
    svg_blip = etree.SubElement(ext, f"{{{SVG_NS}}}svgBlip")
    svg_blip.set(f"{{{R_NS}}}embed", rid)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


REL_FONT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
FACE_KINDS = ("regular", "bold", "italic", "boldItalic")


def _embed_unicodes() -> set[int]:
    keep = set(range(0x20, 0x7F)) | set(range(0xA0, 0x100)) | set(range(0x370, 0x400))
    keep |= set(range(0x2000, 0x2070)) | set(range(0x2190, 0x2300))
    keep |= set(range(0x2460, 0x2500)) | set(range(0x25A0, 0x2600)) | set(range(0x3000, 0x3040))
    keep |= set(range(0x3130, 0x3190))  # Hangul compatibility Jamo
    keep |= set(range(0xAC00, 0xD7A4))  # every Hangul syllable
    keep |= set(range(0xFF00, 0xFF60))
    return keep


def _subset_ttf(path: Path) -> bytes:
    from fontTools import subset
    from fontTools.ttLib import TTFont
    font = TTFont(str(path))
    if "glyf" not in font:
        raise ValueError(f"{path}: not a TrueType-outline font (CFF); PowerPoint cannot embed it")
    if "fvar" in font:
        raise ValueError(f"{path}: variable font; a static instance is needed")
    if font["OS/2"].fsType & 0x0002:
        raise ValueError(f"{path}: fsType forbids embedding")
    options = subset.Options()
    options.layout_features = ["*"]
    options.name_IDs = ["*"]
    options.notdef_outline = True
    options.glyph_names = False
    options.hinting = True
    subsetter = subset.Subsetter(options=options)
    subsetter.populate(unicodes=_embed_unicodes())
    subsetter.subset(font)
    buffer = io.BytesIO()
    font.save(buffer)
    return buffer.getvalue()


def _ttf_to_eot(ttf: bytes) -> bytes:
    """Wrap TrueType data in an EOT 0x00020002 header without compression or XOR: the form
    PowerPoint reads from ppt/fonts/*.fntdata."""
    from fontTools.ttLib import TTFont
    font = TTFont(io.BytesIO(ttf))
    os2, head, name = font["OS/2"], font["head"], font["name"]
    family = name.getDebugName(1) or ""
    style = name.getDebugName(2) or ""
    version = name.getDebugName(5) or ""
    full = name.getDebugName(4) or family
    panose = bytes(getattr(os2.panose, f"b{key}") for key in (
        "FamilyType", "SerifStyle", "Weight", "Proportion", "Contrast",
        "StrokeVariation", "ArmStyle", "LetterForm", "Midline", "XHeight"))
    body = struct.pack("<II", 0x00020002, 0) + panose
    body += struct.pack("<BB", 1, 1 if head.macStyle & 2 else 0)
    body += struct.pack("<IHH", os2.usWeightClass, os2.fsType, 0x504C)
    body += struct.pack("<IIII", os2.ulUnicodeRange1, os2.ulUnicodeRange2, os2.ulUnicodeRange3, os2.ulUnicodeRange4)
    body += struct.pack("<II", getattr(os2, "ulCodePageRange1", 0), getattr(os2, "ulCodePageRange2", 0))
    body += struct.pack("<I", head.checkSumAdjustment) + struct.pack("<IIII", 0, 0, 0, 0)
    for text in (family, style, version, full):
        encoded = text.encode("utf-16-le")
        body += struct.pack("<HH", 0, len(encoded)) + encoded
    body += struct.pack("<HHI", 0, 0, 0)  # padding, root string (empty), root string checksum
    body += struct.pack("<I", 0)  # EUDC code page
    body += struct.pack("<HH", 0, 0)  # padding, signature (empty)
    body += struct.pack("<II", 0, 0)  # EUDC flags, EUDC font size
    return struct.pack("<II", 8 + len(body) + len(ttf), len(ttf)) + body + ttf


def _embedded_face(path: Path) -> bytes:
    """Subsetting a CJK face takes seconds, so the wrapped result is kept per source file
    (path, size, mtime) under the user's cache directory."""
    import hashlib
    stat = path.stat()
    key = hashlib.sha256(f"{path.resolve()}|{stat.st_size}|{stat.st_mtime_ns}".encode("utf-8")).hexdigest()
    cache = Path.home() / ".cache" / "html-mcp-web" / "pptx-fonts" / f"{key}.fntdata"
    if cache.is_file():
        return cache.read_bytes()
    data = _ttf_to_eot(_subset_ttf(path))
    cache.parent.mkdir(parents=True, exist_ok=True)
    cache.write_bytes(data)
    return data


def embed_fonts(pptx_path: Path, family: str, faces: dict[str, Path]) -> None:
    """Add subset TrueType faces to a saved pptx as embedded fonts under `family`."""
    with zipfile.ZipFile(pptx_path) as archive:
        items = {name: archive.read(name) for name in archive.namelist()}
    content_types = items["[Content_Types].xml"].decode("utf-8")
    if 'Extension="fntdata"' not in content_types:
        content_types = content_types.replace(
            "</Types>", '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>')
    rels = items["ppt/_rels/presentation.xml.rels"].decode("utf-8")
    next_id = max((int(value) for value in re.findall(r'Id="rId(\d+)"', rels)), default=0) + 1
    presentation = items["ppt/presentation.xml"].decode("utf-8")
    face_xml = ""
    part_number = 1
    for kind in (kind for kind in FACE_KINDS if kind in faces):
        # A template may carry embedded fonts of its own; their parts keep their names.
        while f"ppt/fonts/font{part_number}.fntdata" in items:
            part_number += 1
        part = f"ppt/fonts/font{part_number}.fntdata"
        items[part] = _embedded_face(faces[kind])
        rid = f"rId{next_id}"
        next_id += 1
        rels = rels.replace(
            "</Relationships>",
            f'<Relationship Id="{rid}" Type="{REL_FONT}" Target="fonts/font{part_number}.fntdata"/></Relationships>')
        face_xml += f'<p:{kind} r:id="{rid}"/>'
    entry = f'<p:embeddedFont><p:font typeface="{family}" pitchFamily="34" charset="-127"/>{face_xml}</p:embeddedFont>'
    if "<p:embeddedFontLst>" in presentation:
        # Keep the template's own entries; an earlier entry for this family gives way.
        presentation = re.sub(
            rf'<p:embeddedFont><p:font typeface="{re.escape(family)}"[^>]*/>.*?</p:embeddedFont>', "", presentation, flags=re.S)
        presentation = presentation.replace("</p:embeddedFontLst>", entry + "</p:embeddedFontLst>", 1)
    else:
        font_list = f"<p:embeddedFontLst>{entry}</p:embeddedFontLst>"
        # Schema order puts embeddedFontLst right after notesSz.
        anchor = re.search(r"<p:notesSz[^>]*/>", presentation)
        if anchor is None:
            raise ValueError("presentation.xml has no notesSz element")
        presentation = presentation[:anchor.end()] + font_list + presentation[anchor.end():]
    root = re.search(r"<p:presentation\b[^>]*>", presentation).group(0)
    if 'embedTrueTypeFonts="1"' not in root:
        presentation = presentation.replace("<p:presentation ", '<p:presentation embedTrueTypeFonts="1" ', 1)
        root = re.search(r"<p:presentation\b[^>]*>", presentation).group(0)
    if "xmlns:r=" not in root:
        presentation = presentation.replace("<p:presentation ", f'<p:presentation xmlns:r="{NS_R}" ', 1)
    items["[Content_Types].xml"] = content_types.encode("utf-8")
    items["ppt/_rels/presentation.xml.rels"] = rels.encode("utf-8")
    items["ppt/presentation.xml"] = presentation.encode("utf-8")
    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", items.pop("[Content_Types].xml"))
        for name, data in items.items():
            archive.writestr(name, data)


def load_pptx_config(skin_dir: Path | None) -> dict[str, Any]:
    if skin_dir is None or not (skin_dir / "skin.json").is_file():
        return {}
    config = json.loads((skin_dir / "skin.json").read_text(encoding="utf-8"))
    return dict(config.get("pptx", {}))


def export_pptx(html_url: str, out_path: Path, project_dir: Path, skin_dir: Path | None) -> dict[str, Any]:
    from marionette_driver.marionette import Marionette
    from pptx import Presentation

    if shutil.which("firefox") is None:
        raise FileNotFoundError("firefox is not installed")
    config = load_pptx_config(skin_dir)
    faces: dict[str, Path] = {}
    if "fonts" in config:
        faces = {kind: skin_dir / config["fonts"][kind] for kind in FACE_KINDS if kind in config["fonts"]}
        for kind, path in faces.items():
            if not path.is_file():
                raise FileNotFoundError(f"pptx font {kind!r}: {path} does not exist")
        font = config["fonts"]["family"]
    else:
        font = config.get("font", DEFAULT_FONT)

    # A port of its own for this firefox. On the default one, an export attached to whatever
    # browser already held 2828: a second export, or one whose predecessor had not finished
    # exiting, drove the wrong session and read another deck's pages.
    with socket.socket() as probe:
        probe.bind(("127.0.0.1", 0))
        marionette_port = int(probe.getsockname()[1])
    profile = tempfile.mkdtemp(prefix="html_mcp_pptx_")
    Path(profile, "user.js").write_text(
        f'user_pref("layout.css.devPixelsPerPx", "{DPR}");\n'
        f'user_pref("marionette.port", {marionette_port});\n'
        'user_pref("browser.shell.checkDefaultBrowser", false);\n', encoding="utf-8")
    proc = subprocess.Popen(
        ["firefox", "-marionette", "-headless", "-no-remote", "-profile", profile,
         "-width", "1400", "-height", "900", "about:blank"],
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    try:
        client = Marionette(host="127.0.0.1", port=marionette_port, startup_timeout=60)
        client.start_session()
        client.navigate(html_url)
        client.execute_script(JS_SETUP)
        for _ in range(100):
            if client.execute_script(JS_READY):
                break
            time.sleep(0.2)
        page_count = client.execute_script(JS_PAGES)
        if page_count == 0:
            raise ValueError("the artifact has no section.page")

        # Every skin renders the same way: the deck's own pages become the slides, so the
        # skin's look comes through its HTML and no skin carries a pptx of its own.
        prs = Presentation()
        prs.slide_width = _px(PAGE_WIDTH_PX)
        prs.slide_height = _px(PAGE_HEIGHT_PX)
        report = []
        for index in range(page_count):
            info = client.execute_script(JS_EXTRACT, script_args=[index])
            if not info["hasBody"]:
                slide = _blank_slide(prs)
                # A full-bleed page's plain text becomes editable text boxes; the rest (CSS
                # decorations, background art, and text carrying a ::before mark) stays in the
                # picture. Hiding the plain text before the shot keeps it from showing twice.
                overlaid = [item for item in info["items"]
                            if item["kind"] == "text" and item.get("plain") and not item.get("pseudo")]
                client.execute_script(JS_HIDE, script_args=[index, [item["i"] for item in overlaid]])
                shot = client.screenshot(element=client.execute_script(
                    "const p = document.querySelectorAll('section.page')[arguments[0]]; p.scrollIntoView(); return p;",
                    script_args=[index]), format="base64")
                _set_slide_background(slide, base64.b64decode(shot))
                for item in overlaid:
                    _add_text(slide, item, font)
                report.append({"page": index + 1, "title": info["title"], "shapes": len(overlaid), "screenshot": True})
                continue
            slide = _blank_slide(prs)
            overlay = [item for item in info["items"] if item.get("body")]
            # An element screenshot (an svg, a KaTeX block) is taken while the element is
            # still visible; the body is then hidden so the chrome-only background does not
            # carry a second copy of what is about to become an editable shape over it.
            for item in overlay:
                if item["kind"] == "shot":
                    element = client.execute_script(JS_ELEMENT, script_args=[index, item["i"]])
                    item["_shot"] = client.screenshot(element=element, format="base64")
            client.execute_script(JS_HIDE, script_args=[index, [item["i"] for item in overlay]])
            shot = client.screenshot(element=client.execute_script(
                "const p = document.querySelectorAll('section.page')[arguments[0]]; p.scrollIntoView(); return p;",
                script_args=[index]), format="base64")
            _set_slide_background(slide, base64.b64decode(shot))
            for item in overlay:
                if item["kind"] == "text":
                    _add_text(slide, item, font)
                elif item["kind"] == "table":
                    _add_table(slide, item, font)
                elif item["kind"] == "panel":
                    _add_panel(slide, item)
                elif item["kind"] == "img":
                    src = item["src"] or ""
                    if src.startswith("data:"):
                        data = base64.b64decode(src.split(",", 1)[1])
                    else:
                        source = (project_dir / src.split("?", 1)[0]).resolve()
                        if not source.is_relative_to(project_dir.resolve()) or not source.is_file():
                            raise FileNotFoundError(f"image {src!r} is not a file inside the project")
                        data = source.read_bytes()
                    fit = item.get("fit")
                    natural = item.get("natural")
                    picture = _add_picture(slide, data, _painted_rect(item["rect"], natural, fit))
                    if fit == "cover":
                        _crop_to_cover(picture, item["rect"], natural)
                else:
                    shot_bytes = base64.b64decode(item["_shot"])
                    if "svg" in item:
                        svg = _normalize_svg(item["svg"].encode("utf-8"), item["rect"][2], item["rect"][3], font)
                        _add_svg_picture(slide, shot_bytes, svg, item["rect"])
                    else:
                        _add_picture(slide, shot_bytes, item["rect"])
            report.append({"page": index + 1, "title": info["title"], "shapes": len(overlay),
                           "screenshot": True, "vector_svgs": sum(1 for item in overlay if "svg" in item)})
        client.delete_session()
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        if faces:
            embed_fonts(out_path, font, faces)
        return {"path": str(out_path), "slides": len(prs.slides), "pages": report,
                "font": font, "embedded_faces": sorted(faces)}
    finally:
        proc.terminate()
        shutil.rmtree(profile, ignore_errors=True)
