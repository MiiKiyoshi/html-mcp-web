import re
import shutil
from pathlib import Path

import pytest

pytest.importorskip("marionette_driver.marionette")
pptx = pytest.importorskip("pptx")

from html_mcp_web.pptx_export import export_pptx  # noqa: E402

REPO = Path(__file__).resolve().parents[1]
NEUTRAL = REPO / "templates" / "neutral-slides"

CONTENT = '''<!doctype html>
<meta charset="utf-8">
<title>Export Deck</title>
<body data-author="Export Author" data-meta="Lab|Today" data-sub="Subtitle line">
<section data-layout="contents" data-title="Contents"><ol><li>Numbers</li></ol></section>
<section data-title="Numbers">
  <p class="lead">Lead sentence with <b>bold</b> and <code>code</code>.</p>
  <ul class="notes"><li>First point<ul><li>Nested point</li></ul></li><li>Second point</li></ul>
  <ol><li>Step one</li><li>Step two</li></ol>
  <table><thead><tr><th>Name</th><th>Value</th></tr></thead>
    <tbody><tr><td>alpha</td><td>1</td></tr><tr><td>beta</td><td>2</td></tr></tbody></table>
  <div class="card"><h3>Card title</h3><p>Card body text.</p></div>
  <p>Math stays a picture: $x^2$</p>
  <img src="fig/box.png" alt="">
</section>
</body>
'''


def deck(tmp_path: Path) -> Path:
    from html_mcp_web.slides import build
    from PIL import Image

    (tmp_path / "fig").mkdir()
    Image.new("RGB", (120, 60), "#3366aa").save(tmp_path / "fig" / "box.png")
    content = tmp_path / "content.html"
    content.write_text(CONTENT, encoding="utf-8")
    output = tmp_path / "slides.html"
    build(content, output, NEUTRAL)
    return output


def shapes_by_kind(slide) -> dict[str, list]:
    kinds: dict[str, list] = {}
    for shape in slide.shapes:
        kinds.setdefault(str(shape.shape_type).split(" ")[0], []).append(shape)
    return kinds


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_export_bakes_chrome_into_background_and_keeps_body_editable(tmp_path: Path) -> None:
    html = deck(tmp_path)
    out = tmp_path / "out" / "deck.pptx"
    result = export_pptx(html.as_uri(), out, tmp_path, None)
    assert result["path"] == str(out)
    assert result["slides"] == 3
    prs = pptx.Presentation(str(out))
    assert (prs.slide_width, prs.slide_height) == (12192000, 6858000)
    # A full-bleed page (cover, contents) is a locked background screenshot the mouse cannot
    # grab, with its words overlaid as editable text boxes so decorations stay in the picture
    # while the text stays real.
    from pptx.oxml.ns import qn as _qn
    for slide in list(prs.slides)[:2]:
        blip = slide._element.find(_qn("p:cSld")).find(_qn("p:bg"))
        assert blip is not None and blip.find(f".//{_qn('a:blip')}") is not None
        assert list(shapes_by_kind(slide)) == ["TEXT_BOX"]
    cover_texts = [s.text_frame.text for s in prs.slides[0].shapes]
    assert "Export Deck" in cover_texts and "Export Author" in cover_texts
    body = prs.slides[2]
    # The chrome (title bar, footer, page number, corner logos) bakes into a locked
    # background the mouse cannot grab; only the body's own blocks stay editable over it.
    body_bg = body._element.find(_qn("p:cSld")).find(_qn("p:bg"))
    assert body_bg is not None and body_bg.find(f".//{_qn('a:blip')}") is not None
    kinds = shapes_by_kind(body)
    texts = [shape.text_frame.text for shape in kinds["TEXT_BOX"]]
    assert "Numbers" in texts   # the bar's title stays editable: decks are retitled per page
    assert "3 / 3" not in texts  # the bar itself and the footer page number baked in
    lead = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("Lead sentence"))
    runs = lead.text_frame.paragraphs[0].runs
    assert [run.text for run in runs] == ["Lead sentence with ", "bold", " and ", "code", "."]
    assert runs[1].font.bold is True
    assert runs[3].font.name == "Consolas"
    assert runs[0].font.name == "Arial"
    bullets = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("First point"))
    assert [p.text for p in bullets.text_frame.paragraphs] == ["First point", "Nested point", "Second point"]
    steps = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("Step one"))
    assert "buAutoNum" in steps.text_frame.paragraphs[0]._p.xml
    assert len(kinds["TABLE"]) == 1
    table = kinds["TABLE"][0].table
    assert table.cell(0, 0).text == "Name" and table.cell(2, 1).text == "2"
    # The table carries no style-drawn grid; only the per-cell borders from the DOM show.
    from pptx.oxml.ns import qn
    assert kinds["TABLE"][0]._element.xpath(".//a:tableStyleId")[0].text == "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"
    # Borders come before the fill in tcPr; PowerPoint drops borders placed after it.
    header_cell = table.cell(0, 0)._tc
    tc_pr = header_cell.find(qn("a:tcPr"))
    order = [child.tag for child in tc_pr]
    fill_index = next(i for i, tag in enumerate(order) if tag.endswith("}noFill") or tag.endswith("}solidFill"))
    assert order.index(qn("a:lnB")) < fill_index
    assert tc_pr.find(qn("a:lnB")).find(qn("a:solidFill")) is not None  # a visible line, not noFill
    # The <code> chip's background becomes a run highlight, and the header's letter-spacing
    # becomes run tracking (spc).
    body_xml = body._element.xml
    assert "a:highlight" in body_xml
    assert re.search(r'spc="\d+"', body_xml)
    assert 'kern="0"' in body_xml  # kerning on, so Latin reads as tight as on the page
    # The card is a panel (its accent rule) followed by its own text; math and the image are pictures.
    assert "Card title" in texts and "Card body text." in texts
    assert len(kinds["PICTURE"]) == 2
    assert kinds["AUTO_SHAPE"]  # the card's accent rule; the chrome bars are baked, not shapes
    # The editable title did not drag its bar back out as a grabbable shape: every panel on
    # the slide belongs to the body, none spans the page width at the top or bottom edge.
    for panel in kinds["AUTO_SHAPE"]:
        spans_page = panel.width >= prs.slide_width - 9525
        assert not spans_page, "a full-width bar came back as a shape"


DEJAVU = Path("/usr/share/fonts/truetype/dejavu")


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
@pytest.mark.skipif(not (DEJAVU / "DejaVuSans.ttf").is_file(), reason="a static TrueType face is required")
def test_export_embeds_the_skin_face(tmp_path: Path) -> None:
    import json
    import re
    import zipfile

    pytest.importorskip("fontTools")
    html = deck(tmp_path)
    skin = tmp_path / "skin"
    (skin / "fonts").mkdir(parents=True)
    shutil.copy(DEJAVU / "DejaVuSans.ttf", skin / "fonts" / "Sans.ttf")
    shutil.copy(DEJAVU / "DejaVuSans-Bold.ttf", skin / "fonts" / "Sans-Bold.ttf")
    (skin / "skin.css").write_text("", encoding="utf-8")
    (skin / "skin.json").write_text(json.dumps({"pptx": {"fonts": {
        "family": "DejaVu Sans", "regular": "fonts/Sans.ttf", "bold": "fonts/Sans-Bold.ttf"}}}), encoding="utf-8")
    out = tmp_path / "deck.pptx"
    result = export_pptx(html.as_uri(), out, tmp_path, skin)
    assert result["font"] == "DejaVu Sans" and result["embedded_faces"] == ["bold", "regular"]
    with zipfile.ZipFile(out) as archive:
        names = archive.namelist()
        assert "ppt/fonts/font1.fntdata" in names and "ppt/fonts/font2.fntdata" in names
        assert archive.read("ppt/fonts/font1.fntdata")[8:12] == b"\x02\x00\x02\x00"  # EOT 0x00020002
        presentation = archive.read("ppt/presentation.xml").decode("utf-8")
        assert 'Extension="fntdata"' in archive.read("[Content_Types].xml").decode("utf-8")
        assert 'embedTrueTypeFonts="1"' in re.search(r"<p:presentation\b[^>]*>", presentation).group(0)
        assert re.search(r'<p:embeddedFont><p:font typeface="DejaVu Sans"[^>]*/><p:regular r:id="rId\d+"/><p:bold r:id="rId\d+"/></p:embeddedFont>', presentation)
    prs = pptx.Presentation(str(out))  # the archive still opens as a presentation
    lead = next(shape for shape in prs.slides[2].shapes if shape.has_text_frame and shape.text_frame.text.startswith("Lead"))
    assert lead.text_frame.paragraphs[0].runs[0].font.name == "DejaVu Sans"


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_export_route_writes_inside_the_project(tmp_path: Path) -> None:
    import json
    import socket
    import urllib.error
    import urllib.request

    import yaml

    from html_mcp_web.config import load_config
    from html_mcp_web.project_server import SharedProjectServer

    deck(tmp_path)
    with socket.socket() as listener:
        listener.bind(("127.0.0.1", 0))
        port = int(listener.getsockname()[1])
    config_path = tmp_path / ".html-mcp-web.yaml"
    config_path.write_text(yaml.safe_dump({
        "artifacts": {"slides": {"label": "Slides", "layout": "slides", "main": "slides.html"}},
        "port": port,
    }, sort_keys=False), encoding="utf-8")
    shared = SharedProjectServer(load_config(config_path))
    try:
        shared.ensure()
        base = f"http://127.0.0.1:{port}/artifacts/slides/export/pptx"

        def post(payload: dict) -> dict:
            request = urllib.request.Request(base, data=json.dumps(payload).encode("utf-8"),
                                             headers={"Content-Type": "application/json"}, method="POST")
            with urllib.request.urlopen(request, timeout=120) as response:
                return json.loads(response.read().decode("utf-8"))

        result = post({"out": "export/deck.pptx"})
        assert result["path"] == str(tmp_path / "export" / "deck.pptx")
        assert (tmp_path / "export" / "deck.pptx").is_file()
        assert [page["screenshot"] for page in result["pages"]] == [True, True, True]
        with pytest.raises(urllib.error.HTTPError) as error:
            post({"out": "../outside.pptx"})
        assert error.value.code == 400
        # The topbar button fetches this: the file is built at export/<artifact>.pptx and served.
        with urllib.request.urlopen(f"http://127.0.0.1:{port}/artifacts/slides/download/pptx", timeout=120) as response:
            assert response.headers["Content-Disposition"] == 'attachment; filename="slides.pptx"'
            body = response.read()
        assert body[:2] == b"PK" and (tmp_path / "export" / "slides.pptx").read_bytes() == body
    finally:
        shared.stop()


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_each_page_screenshots_its_own_element(tmp_path: Path) -> None:
    """Marks restart on every page; a document-wide lookup once shot page 1's math into page 2."""
    from html_mcp_web.slides import build

    content = tmp_path / "content.html"
    content.write_text('''<!doctype html>
<meta charset="utf-8">
<title>Two Math Pages</title>
<body data-author="A" data-meta="B">
<section data-title="First"><p>Short $a$</p></section>
<section data-title="Second"><p>A much longer paragraph carrying $b^2 + c^2$ so its picture is wider than the first</p></section>
</body>
''', encoding="utf-8")
    html = tmp_path / "slides.html"
    build(content, html, NEUTRAL)
    out = tmp_path / "deck.pptx"
    export_pptx(html.as_uri(), out, tmp_path, None)
    prs = pptx.Presentation(str(out))
    first = shapes_by_kind(prs.slides[1])["PICTURE"][0]
    second = shapes_by_kind(prs.slides[2])["PICTURE"][0]
    assert first.image.blob != second.image.blob
    # Both blocks span the body width; the picture is the block shot at 3x.
    assert first.image.size[0] == second.image.size[0] == round(first.width / 9525 * 3)


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_inline_svg_is_embedded_as_vector_math_stays_raster(tmp_path: Path) -> None:
    import re
    import zipfile

    from html_mcp_web.slides import build

    content = tmp_path / "content.html"
    content.write_text('''<!doctype html>
<meta charset="utf-8">
<title>Vector</title>
<body data-author="A" data-meta="B">
<section data-title="Diagram">
  <svg viewBox="0 0 200 100" xmlns="http://www.w3.org/2000/svg"><rect x="10" y="10" width="180" height="80" fill="none" stroke="#333"/></svg>
</section>
<section data-title="Formula"><p>Only math here: $x^2 + y^2$</p></section>
</body>
''', encoding="utf-8")
    html = tmp_path / "slides.html"
    build(content, html, NEUTRAL)
    out = tmp_path / "deck.pptx"
    result = export_pptx(html.as_uri(), out, tmp_path, None)
    diagram = next(p for p in result["pages"] if p["title"] == "Diagram")
    formula = next(p for p in result["pages"] if p["title"] == "Formula")
    assert diagram["vector_svgs"] == 1 and formula["vector_svgs"] == 0
    with zipfile.ZipFile(out) as archive:
        svg_parts = [n for n in archive.namelist() if n.endswith(".svg")]
        assert len(svg_parts) == 1
        # The embedded svg's viewBox aspect is normalized to the picture frame so PowerPoint,
        # which fills the frame, does not stretch the drawing.
        svg_text = archive.read(svg_parts[0]).decode("utf-8")
        vb = [float(v) for v in re.search(r'viewBox="([^"]+)"', svg_text).group(1).split()]
        pic = next(sh for sh in pptx.Presentation(str(out)).slides[1].shapes if sh.shape_type is not None and "PICTURE" in str(sh.shape_type))
        assert abs(vb[2] / vb[3] - pic.width / pic.height) < 0.01
        # The deck font is named on the svg root so PowerPoint keeps the browser's metrics
        # for the labels instead of substituting a wider face.
        assert re.search(r'<svg[^>]*font-family="[^"]+"', svg_text)
        assert "image/svg+xml" in archive.read("[Content_Types].xml").decode("utf-8")
        # The diagram slide's blip carries an svgBlip pointing at the svg part; math does not.
        diagram_xml = archive.read("ppt/slides/slide2.xml").decode("utf-8")
        formula_xml = archive.read("ppt/slides/slide3.xml").decode("utf-8")
        assert "svgBlip" in diagram_xml and "svgBlip" not in formula_xml
        svg_rid = re.search(r'svgBlip[^>]*r:embed="(rId\d+)"', diagram_xml).group(1)
        rels = archive.read("ppt/slides/_rels/slide2.xml.rels").decode("utf-8")
        target = re.search(rf'Id="{svg_rid}"[^>]*Target="([^"]+)"', rels).group(1)
        assert target.endswith(".svg")
        # both slides still carry a raster fallback picture
        assert "<p:pic>" in diagram_xml and "<p:pic>" in formula_xml
    pptx.Presentation(str(out))  # the package still opens


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_inline_flow_images_each_become_a_picture(tmp_path: Path) -> None:
    """A row of <img> in inline flow (no block wrapper) reads as a text-only leaf; each image
    must still export as its own picture instead of being dropped with the empty text."""
    from html_mcp_web.slides import build
    from PIL import Image

    (tmp_path / "fig").mkdir()
    for name in ("a", "b", "c", "d"):
        Image.new("RGB", (80, 80), "#3366aa").save(tmp_path / "fig" / f"{name}.png")
    content = tmp_path / "content.html"
    content.write_text('''<!doctype html>
<meta charset="utf-8">
<title>Inline Images</title>
<body data-author="A" data-meta="B">
<section data-title="Row">
  <div><img src="fig/a.png" style="width:20%"><img src="fig/b.png" style="width:20%"><img src="fig/c.png" style="width:20%"><img src="fig/d.png" style="width:20%"></div>
</section>
</body>
''', encoding="utf-8")
    html = tmp_path / "slides.html"
    build(content, html, NEUTRAL)
    out = tmp_path / "deck.pptx"
    export_pptx(html.as_uri(), out, tmp_path, None)
    prs = pptx.Presentation(str(out))
    assert len(shapes_by_kind(prs.slides[1]).get("PICTURE", [])) == 4


@pytest.mark.skipif(shutil.which("firefox") is None, reason="Firefox is required")
def test_the_page_is_carried_over_as_it_is_laid_out(tmp_path: Path) -> None:
    """Four things the export used to lose: the line breaks of a code block, the gap under
    a bullet, the proportions of a picture fitted in a box of another shape, and a drawing
    whose only fault was sitting inline, which came out as a line of its own labels."""
    from html_mcp_web.slides import build
    from PIL import Image

    (tmp_path / "fig").mkdir()
    Image.new("RGB", (120, 60), "#3366aa").save(tmp_path / "fig" / "wide.png")
    content = tmp_path / "content.html"
    content.write_text('''<!doctype html>
<meta charset="utf-8">
<title>Fidelity</title>
<body data-author="A" data-meta="B">
<section data-title="Carried Over">
  <pre><code>first (line) {
  second : 50;
  third  : 90;
}</code></pre>
  <ul class="notes"><li>one</li><li>two</li><li>three</li></ul>
  <img src="fig/wide.png" style="width:300px;height:100px;object-fit:contain">
  <div><svg viewBox="0 0 200 100" width="200" height="100"><rect x="5" y="5" width="190" height="90" fill="none" stroke="#333"/><text x="10" y="60" font-size="12">axis label</text></svg></div>
  <p>step <span style="display:inline-block;min-width:1.25em;height:1.25em;text-align:center;border-radius:50%;background:#243943;color:#fff;font-size:13px">2</span> then <code>chip</code></p>
</section>
</body>
''', encoding="utf-8")
    html = tmp_path / "slides.html"
    build(content, html, NEUTRAL)
    out = tmp_path / "deck.pptx"
    export_pptx(html.as_uri(), out, tmp_path, None)
    slide = pptx.Presentation(str(out)).slides[1]
    kinds = shapes_by_kind(slide)

    listing = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("first (line)"))
    assert listing.text_frame.text.count("\v") == 3  # four lines, so three breaks
    assert "  second : 50;" in listing.text_frame.text  # and the indentation with them

    bullets = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("one"))
    gaps = [para.space_before for para in bullets.text_frame.paragraphs]
    assert gaps[0].pt == 0
    assert all(gap.pt > 0 for gap in gaps[1:])  # ul.notes spaces its items with margin-bottom

    # The image is 2:1 in a 3:1 box, so the page draws it 200x100 and leaves a band at each
    # side; the picture carries the drawn size, not the box.
    picture = next(shape for shape in kinds["PICTURE"]
                   if shape.image.blob[:4] == b"\x89PNG" and shape.image.size == (120, 60))
    assert round(picture.width / 9525) == 200 and round(picture.height / 9525) == 100

    drawing = [shape for shape in kinds["PICTURE"] if shape.image.size != (120, 60)]
    assert drawing, "the inline svg is a picture of its own"
    assert not any("axis label" in shape.text_frame.text for shape in kinds["TEXT_BOX"])

    # A numbered marker is a disc with a pale digit on it. A run carries no shape, and a
    # highlight is a rectangle that some readers carry past a line break, so the character
    # that means this is used: it sits in the line and cannot drift from it.
    from pptx.oxml.ns import qn
    line = next(shape for shape in kinds["TEXT_BOX"] if shape.text_frame.text.startswith("step"))
    runs = line.text_frame.paragraphs[0].runs
    badge = next(run for run in runs if run.text == "\u2777")  # the disc for 2
    assert str(badge.font.color.rgb) == "243943"
    assert badge._r.find(qn("a:rPr")).find(qn("a:highlight")) is None
    chip = next(run for run in runs if run.text == "chip")
    assert chip._r.find(qn("a:rPr")).find(qn("a:highlight")) is not None  # a code chip still is one
